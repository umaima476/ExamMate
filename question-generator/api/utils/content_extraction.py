import fitz
from pptx import Presentation
import io

def extract_text(file_bytes, filename):
    text = ""
    if filename.endswith(".pdf"):
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page in doc:
            text += page.get_text()
    
    elif filename.endswith(".pptx") or filename.endswith(".ppt"):
        prs = Presentation(io.BytesIO(file_bytes))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
    return text